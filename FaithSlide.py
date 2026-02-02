from docx import Document
from pptx import Presentation
from pptx.util import Pt
import copy
import os
from threading import Thread
import logging
from time import sleep
import sys
from random import uniform
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import re
import requests

if getattr(sys, 'frozen', False):
    base_path = sys._MEIPASS
    exe_dir = os.path.dirname(sys.executable)
    log_path = os.path.join(exe_dir, "FaithSlide.log")
else:
    base_path = os.path.dirname(os.path.abspath(__file__))
    log_path = os.path.join(base_path, "FaithSlide.log")

logging.basicConfig(
    filename=log_path,
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    encoding="utf-8"
)
template_ppt_file = os.path.join(base_path, "template.pptx")

prs = None
book_var = None
chapter_var = None
verse_var = None
text_box = None
progress_bar = None
url = "https://bible.fhl.net/json/qb.php"
# 簡稱 -> 全名
abbr_to_full = {
    "創": "創世記",
    "出": "出埃及記",
    "利": "利未記",
    "民": "民數記",
    "申": "申命記",
    "書": "約書亞記",
    "士": "士師記",
    "得": "路得記",
    "撒上": "撒母耳記上",
    "撒下": "撒母耳記下",
    "王上": "列王紀上",
    "王下": "列王紀下",
    "代上": "歷代志上",
    "代下": "歷代志下",
    "拉": "以斯拉記",
    "尼": "尼希米記",
    "斯": "以斯帖記",
    "伯": "約伯記",
    "詩": "詩篇",
    "箴": "箴言",
    "傳": "傳道書",
    "歌": "雅歌",
    "賽": "以賽亞書",
    "耶": "耶利米書",
    "哀": "耶利米哀歌",
    "結": "以西結書",
    "但": "但以理書",
    "何": "何西阿書",
    "珥": "約珥書",
    "摩": "阿摩司書",
    "俄": "俄巴底亞書",
    "拿": "約拿書",
    "彌": "彌迦書",
    "鴻": "那鴻書",  # 小先知書，部分版本略有不同
    "哈": "哈巴谷書",
    "番": "西番雅書",
    "該": "哈該書",
    "瑪": "瑪拉基書",
    "亞": "撒迦利亞書",
    "太": "馬太福音",
    "可": "馬可福音",
    "路": "路加福音",
    "約": "約翰福音",
    "徒": "使徒行傳",
    "羅": "羅馬書",
    "林前": "哥林多前書",
    "林後": "哥林多後書",
    "加": "加拉太書",
    "弗": "以弗所書",
    "腓": "腓立比書",
    "西": "歌羅西書",
    "帖前": "帖撒羅尼迦前書",
    "帖後": "帖撒羅尼迦後書",
    "提前": "提摩太前書",
    "提後": "提摩太後書",
    "多": "提多書",
    "門": "腓利門書",
    "來": "希伯來書",
    "雅": "雅各書",
    "彼前": "彼得前書",
    "彼後": "彼得後書",
    "約壹": "約翰一書",
    "約貳": "約翰二書",
    "約參": "約翰三書",
    "約一": "約翰一書",
    "約二": "約翰二書",
    "約三": "約翰三書",
    "猶": "猶大書",
    "啟": "啟示錄"
}
# 全名 -> 簡稱
full_to_abbr = {v: k for k, v in abbr_to_full.items()}
chinese_number = ["零", "一", "二", "三", "四", "五", "六", "七", "八", "九", "十"]
number = ["1", "2", "3", "4", "5", "6", "7", "8", "9", "0"]
# 旧约书卷列表
search_page = False
# all_book = ["創", "出", "利", "民", "申", "書", "士", "得", "撒上", "撒下", "王上", "王下", "代上", "代下", "拉", "尼", "斯", "伯", "詩", "箴", "傳", "歌", "賽", "耶", "哀", "結", "但", "何", "珥", "摩", "俄", "拿", "彌", "鴻", "哈", "番", "該", "瑪", "亞", "太", "可", "路", "約", "徒", "羅", "林前", "林後", "加", "弗", "腓", "西", "帖前", "帖後", "提前", "提後", "多", "門", "來", "雅", "彼前", "彼後", "約壹", "約貳", "約參", "猶", "啟"]
books = "創|出|利|民|申|書|士|得|撒上|撒下|王上|王下|代上|代下|拉|尼|斯|伯|詩|箴|傳|歌|賽|耶|哀|結|但|何|珥|摩|俄|拿|彌|鴻|哈|番|該|瑪|亞|太|可|路|約|徒|羅|林前|林後|加|弗|腓|西|帖前|帖後|提前|提後|多|門|來|雅|各|彼前|彼後|約壹|約貳|約參|猶|啟|約一|約二|約三"
main_book = ""

#爬蟲抓經文
def get_verses(book_abbr, chapter):
    try:
        if book_abbr in ["各", "約壹", "約貳", "約參"]:
            index = ["各", "約壹", "約貳", "約參"].index(book_abbr)
            book_abbr = ["雅", "約一", "約二", "約三"][index]

        params = {
            "chineses": book_abbr,   # 書卷
            "chap": chapter,        # 章
            "version": "unv",   # 強制指定為和合本 (Union Version)
            "strong": "0"       # 明確要求不要 Strong Number
        }

        # 執行請求
        response = requests.get(url, params=params, timeout=5)
        
        if response.status_code == 200:
            data = response.json() # 這行最關鍵！直接把結果變字典
            verses = []

            if data.get('status') == 'success':
                for record in data["record"]:
                    verses.append(f"{record['sec']}. {record['bible_text']}")
                return verses
            else:
                logging.warning(f"API 狀態錯誤: {data.get('status')}")
        return []
    except Exception as e:
        logging.error(f"get_verses 發生異常: {e}")
        return []
#PPT 複製投影片
def duplicate_slide(prs:Presentation, index):
    try:
        template_slide = prs.slides[index]
        new_slide = prs.slides.add_slide(template_slide.slide_layout)
        for shape in list(new_slide.shapes):
            if shape.is_placeholder:
                sp = shape
                new_slide.shapes._spTree.remove(sp._element)
        for shape in template_slide.shapes:
            # if not shape.is_placeholder:
                el = shape.element
                new_el = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        return new_slide
    except Exception as e:
        logging.warning(f"duplicate_slide {e}")
        messagebox.showwarning("錯誤", "複製PPT時錯誤")
#PPT 刪除投影片
def remove_slide(prs:Presentation, index:int) -> None:
    try:
        xml_slides = prs.slides._sldIdLst
        slide = list(xml_slides)
        xml_slides.remove(slide[index])
    except Exception as e:
        logging.warning(f"remove_slide {e}")
        messagebox.showwarning("錯誤", "刪除PPT時錯誤")
def get_weighted_length(text: str) -> float:
    length = 0.0
    for char in text:
        if re.match(r'[\u4e00-\u9fff]|["，。？！：「」；]', char):
            length += 1.0  # Chinese character
        else:
            length += 0.5  # Non-Chinese character
    # print(length)
    return length

def calculate_font_size(text: str, current_base_size=72) -> Pt:
    
    # 參數設定
    MAX_WIDTH_UNITS = 854 # 每行最大中文字符數（需實驗調整）
    MAX_HIGHT_UNITS = 372  # 每頁最大行數（需實驗調整）
    MIN_FONT_SIZE = 20   # 最小可讀字體大小
    
    total_length = get_weighted_length(text)
    
    # 1. 計算縮放因子 (Scaling Factor)
    # 如果長度超過最大單行容量，則需要縮放
    
    left_size = MIN_FONT_SIZE
    right_size = current_base_size
    # 二分法微調字體大小
    while right_size - left_size > 0.01:
        mid_size = (left_size + right_size) / 2
        width_units = int(MAX_WIDTH_UNITS / mid_size)
        hight_units = int(MAX_HIGHT_UNITS / mid_size)
        if total_length > width_units * hight_units:
            right_size = mid_size
        else:
            left_size = mid_size
        new_size = left_size
        # print(f"mid_size: {mid_size}, width_units: {width_units}, hight_units: {hight_units}, total_length: {width_units * hight_units}")



    # 3. 施加約束（防止字體過小或超過基數）
    if new_size < MIN_FONT_SIZE:
        final_size = MIN_FONT_SIZE
    elif new_size > current_base_size:
        final_size = current_base_size
    else:
        final_size = new_size
    # print(final_size)

    return Pt(final_size) # 必須返回 pptx.util.Pt 對象
#PPT 經文投影片
def verses_PPT(title:str, verses:str):
    try:
        if "." not in verses:
            logging.warning(f"經文格式錯誤，無法製作投影片: title: {title} verse: {verses}")
            return
        num = verses.split(".")[0] + "."
        out_verses = verses.split(".")[1]

        if len(num) == 2:
            new_slide = duplicate_slide(prs, 5)
        else:
            new_slide = duplicate_slide(prs, 0)

        text_frame = new_slide.shapes[0].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        if "詩篇" in title:
            title = title.replace("章", "篇")
        p.runs[0].text = title
        for i in range(1, 3):
            try:
                p.runs[i].text = ""
            except:
                break

        text_frame = new_slide.shapes[1].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        
        p.runs[0].text = num
        out_verses = out_verses.replace("　", " ")    
        p.runs[1].text = out_verses

        text_size = calculate_font_size(out_verses)
        p.runs[0].font.size = text_size
        p.runs[1].font.size = text_size
    except Exception as e:
        logging.warning(f"verses_PPT {e}")
        messagebox.showwarning("錯誤", "製作經文PPT時錯誤")
#PPT 主標題投影片
def main_title_PPT(title):
    try:
        new_slide = duplicate_slide(prs, 1)

        text_frame = new_slide.shapes[1].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        p.runs[0].text = title
        new_slide = duplicate_slide(prs, 2)
    except Exception as e:
        logging.warning(f"main_title_PPT {e}")
        messagebox.showwarning("錯誤", "製作主標題PPT時錯誤")
#PPT 大標題投影片
def major_heading_PPT(major):
    try:
        new_slide = duplicate_slide(prs, 3)

        text_frame = new_slide.shapes[0].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        p.runs[0].text = major.split("、")[0] + "、"
        p.runs[1].text = major.split("、")[1]
    except Exception as e:
        logging.warning(f"major_heading_PPT {e}")
        messagebox.showwarning("錯誤", "製作大標題PPT時錯誤")
#PPT 中標題投影片
def medium_hearding_PPT(major, medium, medium_list):
    try:
        new_slide = duplicate_slide(prs, 4)

        text_frame = new_slide.shapes[0].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        p.runs[0].text = major.split("、")[0] + "、"
        p.runs[1].text = major.split("、")[1]

        text_frame = new_slide.shapes[1].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        t = 0
        for m in medium_list:
            p.runs[2*t].text = m.split(".")[0] + "."
            p.runs[2*t+1].text = m.split(".")[1].replace("：", "") + "\n"
            if m == medium:
                break
            t += 1
            p.add_run()
            p.add_run()
    except Exception as e:
        logging.warning(f"medium_hearding_PPT {e}")
        messagebox.showwarning("錯誤", "製作中標題PPT時錯誤")
#PPT 小標題投影片
def minor_heading_PPT(major, medium, minor, minor_list):
    try:
        new_slide = duplicate_slide(prs, 4)

        text_frame = new_slide.shapes[0].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()
        p.runs[0].text = major.split("、")[0] + "、"
        p.runs[1].text = major.split("、")[1]

        text_frame = new_slide.shapes[1].text_frame
        p = text_frame.paragraphs[0]
        if not p.runs:
            p.add_run()

        p.runs[0].text = medium.split(".")[0] + "."
        p.runs[1].text = medium.split(".")[1].replace("：", "") + "\n"

        t = 1
        for m in minor_list:
            p.add_run()
            p.add_run()
            p.runs[2*t].text = "(" + m.split(")")[0] + ")"
            p.runs[2*t+1].text = m.split(")")[1].replace("：", "") + "\n"
            if m == minor:
                break
            t += 1
    except Exception as e:
        logging.warning(f"minor_heading_PPT {e}")
        messagebox.showwarning("錯誤", "製作小標題PPT時錯誤")
#PPT 經文章節轉中文
def num_to_chinese(title, chapter_and_verse: str) -> str:
    try:
        chapter = chapter_and_verse.split(":")[0]
        chinese_chapter = ""
        # print(chapter)
        if len(chapter) == 3:
            chinese_chapter += f"{chinese_number[int(chapter[0])]}百"
            chinese_chapter += f"{chinese_number[int(chapter[1])]}"
            if chinese_chapter[-1] != "零":
                chinese_chapter += "十"
            chinese_chapter += f"{chinese_number[int(chapter[2])]}"
        elif len(chapter) == 2:
            chinese_chapter += f"{chinese_number[int(chapter[0])]}"
            if chinese_chapter == "一":
                chinese_chapter = ""
            if chapter[1] == "0":
                chinese_chapter += "十"
            else:
                chinese_chapter += f"十{chinese_number[int(chapter[1])]}"
        elif len(chapter) == 1:
            chinese_chapter = f"{chinese_number[int(chapter[0])]}"

        title += f"{chinese_chapter}章"
        return title
    except Exception as e:
        logging.warning(f"num_to_chinese {e}")
        messagebox.showwarning("錯誤", "阿拉伯數字轉中文時錯誤")
#PPT 經文節數分析
def analyze_paragraph(title, verse_analyze, verses):
    try:
        start = int(verse_analyze.split("-")[0])-1
        try:
            end = int(verse_analyze.split("-")[1].replace(",",""))
        except:
            end = start + 1
        for v in range(start, end):
            verse = verses[v].replace(" ", "")
            verses_PPT(title, verse)
            logging.info(f"{title} {verse}")
    except Exception as e:
        logging.warning(f"analyze_paragraph {e}")
        messagebox.showwarning("錯誤", "經文節數分析時錯誤")
#PPT 經文章節處理
def process_reference_block(chapter_and_verse, book):
    try:
        scrape_verses = get_verses(book, chapter_and_verse.split(":")[0])
        if scrape_verses == []:
            logging.warning(f"尚未取得經文，跳過，{book} {chapter_and_verse}")
            return
        
        title = f"{abbr_to_full[book]}"
        title = num_to_chinese(title, chapter_and_verse)
        chapter_and_verse = chapter_and_verse.replace("，", "")
        verse = chapter_and_verse.split(':')[1]
        if "," in verse:
            verse = verse.split(",")

        if isinstance(verse, list):
            # print(verse, "is verse list")
            for v in verse:
                if v:
                    analyze_paragraph(title+f"{v}節", v, scrape_verses)
        else:
            title += f"{verse}節"
            analyze_paragraph(title, verse, scrape_verses)
    except Exception as e:
        logging.warning(f"process_reference_block {e}")
        messagebox.showwarning("錯誤", "經文章節處理時錯誤")
#PPT 經文書卷解析
def parse_bible_reference(bible):
    try:
        # print(bible)
        book = ""
        chapter_and_verse = ""
        for char in bible:
            # print(char)
            if char[0] in number:
                chapter_and_verse += char
                if book == "":
                    book = main_book                
                
                if chapter_and_verse.count(":") > 1:
                    cut_time = 0
                    new_chapter_and_verse = []
                    text = ""
                    for t in chapter_and_verse:
                        text += t
                        if t == ":":
                            cut_time += 1
                        if cut_time > 1:
                            verse = "".join(text.split(",")[0:-1])
                            new_chapter_and_verse.append(verse)
                            text = text.replace(verse, "")[1::]
                            cut_time -= 1
                    new_chapter_and_verse.append(text)
                    chapter_and_verse = new_chapter_and_verse
                if isinstance(chapter_and_verse, list):
                    for cav in chapter_and_verse:
                        process_reference_block(cav, book)
                else:
                    process_reference_block(chapter_and_verse, book)                

                book = ""
                chapter_and_verse = ""
            else:
                book = char
    except Exception as e:
        logging.warning(f"parse_bible_reference {e}")
        messagebox.showwarning("錯誤", "經文書卷解析時錯誤")
#PPT 段落處理
def paragraph_PPT(heading, verses):
    try:
        if heading["minor"]:
            heading_livel = 3
        elif heading["medium"]:
            heading_livel = 2
        else:
            heading_livel = 1
        logging.info(verses)
        if heading_livel == 1:
            major_heading_PPT(heading["major"])
            parse_bible_reference(verses[0])
        elif heading_livel == 2:
            major_heading_PPT(heading["major"])
            parse_bible_reference(verses[0])
            for medium in heading["medium"]:
                medium_hearding_PPT(heading["major"], medium, heading["medium"])
                parse_bible_reference(verses[1][medium])
        elif heading_livel == 3: #確認模板
            major_heading_PPT(heading["major"])
            parse_bible_reference(verses[0])
            for medium in heading["medium"]:
                medium_hearding_PPT(heading["major"], medium, heading["medium"])
                parse_bible_reference(verses[1][medium])
                if medium in heading["minor"].keys():
                    for minor in heading["minor"][medium]:
                        minor_heading_PPT(heading["major"], medium, minor, heading["minor"][medium])
                        parse_bible_reference(verses[2][minor])
    except Exception as e:
        logging.warning(f"paragraph_PPT {e}")
        messagebox.showwarning("錯誤", "段落處理時錯誤")
#關閉驅動程式
def close_driver():
    try:
        root.destroy()
    except Exception as e:
        logging.warning(f"close_driver {e}")
        messagebox.showwarning("錯誤", "關閉程式時發生異常，請稍後再試")
#分析word
def analyze_word(text):
    for word in ["證道", "傳道", "牧師", "吳佩倫", "錢致榮"]:
        text = text.replace(word, "")
    parts = re.split(r'(?=[一二三四五六七八九十]+\s*、|\d+\s*\.|\d+\s*\))', text)
    sermon = []
    for p in parts:
        # print(p)
        title = p
        del_matches = re.findall(rf"(?<![\u4e00-\u9fff])(?:{books})+\s*\d+\s*:\s*(?:\d+(?:-\d+)?)+(?:,\s*\d+(?::\s*\d+)*(?:-\s*\d+)?)*|\d+\s*:(?:\s*\d+(?:-\s*\d+)?)+(?:,\s*\d+(?::\s*\d+)*(?:-\s*\d+)?)*", p)
        matches = re.findall(rf"(?<![\u4e00-\u9fff])(?:{books})+(?![\u4e00-\u9fff])|\d+\s*:(?:(?:\s*\d+\s*)(?:-\s*(?:\s*\d+\s*))?)+(?:\s*,\s*(?:\d+\s*)*(?::\s*\d+)*(?:-\s*(?:\d+\s*)*)?)*", p)
        cleaned_matches = [m.replace(" ", "") for m in matches]
        if del_matches:
            title = title.split(del_matches[0])[0].replace(" ", "")

        sermon.append(title.strip())
        sermon.extend(cleaned_matches)
    return sermon
#製作ppt
def produce_the_slides():
    # messagebox.showwarning("開始製作投影片，請稍候...")
    global main_book, prs
    try:
        prs = Presentation(template_ppt_file)
        # print(log_path, "為日誌檔案位置")
        # 請改成你的 Word 路徑
        wordfile_path = word_path_var.get()
        # print(wordfile_path)
        doc = Document(wordfile_path)
        ReadTheBible = []
        text_box = []
        # 逐個表格抓文字
        for t_idx, table in enumerate(doc.tables):
            # print(f"=== 表格 {t_idx+1} ===")
            for r_idx, row in enumerate(table.rows):
                tatil = row.cells[0].text.strip()
                # 只印出有內容的列
                if tatil == "讀經":
                    row_texts = [cell.text.strip() for cell in row.cells]
                    ReadTheBible = row_texts[1].split("\n")
                elif tatil == "證道": 
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                text = run.text.strip()
                                if run.bold and text:
                                    # print(text)
                                    text_box.append(text)
        sermon = analyze_word(" ".join(text_box))

    except Exception as e:
        logging.warning(f"Analyze_and_produce_the_slides {e}")
        messagebox.showwarning("錯誤", "分析word時錯誤")
        return
    
    try:
        update_progress(1, 10, "分析word完畢")

        if not ReadTheBible:
            logging.warning("讀經抓取失敗")
        else:
            logging.info("讀經:")
            main_verses = ReadTheBible[0]
            
            if "，" in main_verses:
                main_book = ""
                for text in main_verses:
                    if text in chinese_number:
                        break
                    main_book += text
                main_verses = main_verses.replace("，", " " + main_book).split()
            del ReadTheBible[0]
            for verses_index in range(0, len(ReadTheBible)):
                ReadTheBible[verses_index] = ReadTheBible[verses_index].replace("[", "").replace("]", ".")
            verses_index = 0

            logging.info(ReadTheBible)
            if not isinstance(main_verses, list):
                for verses in ReadTheBible:
                    logging.info(f"{main_verses}, {verses}")
                    if "." not in verses:
                        main_verses = verses
                        continue
                    verses_PPT(main_verses, verses)
            else:
                for verses in main_verses:
                    first_num = 0
                    second_num = 0
                    first_end = False
                    for text in verses:
                        if text in number:
                            if not first_end:
                                first_num *= 10
                                first_num += int(text)
                            else:
                                second_num *= 10
                                second_num += int(text)
                        elif text == "-":
                            first_end = True
                    # print(verses, first_num, second_num)
                    if second_num == 0:
                        second_num = first_num
                    for i in range(first_num, second_num+1):
                        logging.info(f"{verses}, {ReadTheBible[verses_index]}")
                        verses_PPT(verses, ReadTheBible[verses_index])

                        verses_index += 1
        update_progress(2, 10, "讀經製作完畢")
    except Exception as e:
        logging.warning(f"Analyze_and_produce_the_slides {e}")
        messagebox.showwarning("錯誤", "製作讀經PPT時錯誤")
        return
    
    try:
        for book in full_to_abbr.keys():
            if isinstance(main_verses, list):
                if book in main_verses[0]:
                    main_book = full_to_abbr[book]
                    break
            else:
                if book in main_verses:
                    main_book = full_to_abbr[book]
                    break
        update_progress(2.5, 10, "主經文分析完畢")
        logging.info(f"main book {main_book}")
    except Exception as e:
        logging.warning(f"Analyze_and_produce_the_slides {e}")
        messagebox.showwarning("錯誤", "分析主經文時錯誤")
        return
    
    try:
        if not sermon:
            logging.warning("證道抓取失敗")
        else:
            logging.info(f"證道:{sermon}")
            # print(f"證道:{sermon}")
            make_main_title = False
            heading = {"major": "", "medium": [], "minor": {}}
            verses = [[], {}, {}]  # 大標題，主標題，副標題 經文
            heading_livel = 0
            for step, text in zip(range(1, len(sermon)+1), sermon):
                update_progress(2.5+step*7/len(sermon), 10, heading["major"])
                if not make_main_title: # 大標題
                    main_title_PPT(text)
                    make_main_title = True
                else:
                    if "、" in text: # 主標題
                        if heading_livel != 0:# 已有完整段落，製作PPT
                            logging.info(f"{heading}, {verses}")
                            # print(heading, "\n", verses, "complete paragraph")
                            paragraph_PPT(heading, verses)
                            heading = {"major": "", "medium": [], "minor": {}}
                            verses = [[], {}, {}]  # 大標題，主標題，副標題 經文

                        heading_livel = 1
                        heading["major"] = text
                    elif "." in text: # 副標題
                        if heading["major"] == "":
                            logging.info("副標題出現於主標題之前，格式錯誤")
                        else:  
                            heading["medium"].append(text)
                            
                            subtitle = True
                            heading_livel = 2

                    elif ")" in text:  # 小標題，待測試
                        heading_livel = 3
                        if len(heading["medium"]) == 0:
                            logging.info("小標題出現於副標題之前，格式錯誤")
                        if heading["medium"][-1] not in heading["minor"].keys():
                            heading["minor"][heading["medium"][-1]] = []
                        heading["minor"][heading["medium"][-1]].append(text)
                        minor_title = True

                    else:
                        # print(text, "is verse")
                        is_verse = False
                        for t in text:
                            if t in number:
                                is_verse = True
                                break
                        else:
                            if text in abbr_to_full.keys():
                                is_verse = True
                        if is_verse:
                            if heading_livel == 1:
                                verses[0].append(text)
                            elif heading_livel == 2:
                                if heading["medium"][-1] not in verses[1].keys():
                                    verses[1][heading["medium"][-1]] = []
                                verses[1][heading["medium"][-1]].append(text)
                            else:
                                last_medium = heading["medium"][-1]
                                if heading["minor"][last_medium][-1] not in verses[2].keys():
                                    verses[2][heading["minor"][last_medium][-1]] = []
                                verses[2][heading["minor"][last_medium][-1]].append(text)
                                # print("小標題經文待測試")

            logging.info(f"{heading}, {verses}")
            paragraph_PPT(heading, verses)
            update_progress(10, 10, "製作完畢")
    except Exception as e:
        logging.warning(f"Analyze_and_produce_the_slides {e}")
        messagebox.showwarning("錯誤", "製作證道PPT時錯誤")
        return
        # print(heading, "\n", verses, "final paragraph")
        
                    
    # 刪除範本投影片                     
    for _ in range(6):
        remove_slide(prs,0)
    try:
        save_path = ppt_save_var.get()
        prs.save(save_path)
        logging.info("製作完成")
    except Exception as e:
        logging.warning(f"Analyze_and_produce_the_slides {e}")
        messagebox.showwarning("錯誤", "PPT存檔錯誤")
        return
    messagebox.showinfo("", "製作完成")
#以另一線程製作PPT
def Start_produce():
    Thread(target=produce_the_slides, daemon=True).start()
#清空UI介面
def clear_frame(frame_to_clear):
    try:
        for widget in frame_to_clear.winfo_children():
            widget.destroy()
    except Exception as e:
        logging.warning(f"clear_frame {e}")
        messagebox.showwarning("錯誤", "清空UI時錯誤")
#經文搜尋工具
def run_search():
    try:
        book_abbr = book_var.get()
        chapter = chapter_var.get()
        verse = verse_var.get()

        if not book_abbr or not chapter:
            messagebox.showwarning("輸入錯誤", "請選擇書卷與章節")
            logging.error("輸入錯誤", "請選擇書卷與章節")
            return
        if book_abbr not in abbr_to_full.keys():
            book_abbr = full_to_abbr.get(book_abbr, "")
            if not book_abbr:
                messagebox.showwarning("輸入錯誤", "書卷名稱錯誤")
                logging.error("輸入錯誤", "書卷名稱錯誤")
                return

        verses = get_verses(book_abbr, chapter)

        text_box.delete(1.0, tk.END)
        if verse:
            if "-" in verse:
                start, end = map(int, verse.split("-"))
                for v in range(start, end + 1):
                    if 1 <= v <= len(verses):
                        text_box.insert(tk.END, verses[v - 1] + "\n")
                    else:
                        messagebox.showwarning("輸入錯誤", "節數錯誤")
                        logging.warning("抓取經文失敗: 節數超出範圍")
                        break
                logging.info(f"成功抓取 {abbr_to_full[book_abbr]} 第 {chapter} 章 {start}-{end} 節")
            else:
                v= int(verse)
                if 1 <= v <= len(verses):
                    text_box.insert(tk.END, verses[v-1] + "\n")
                    logging.info(f"成功抓取 {abbr_to_full[book_abbr]} 第 {chapter} 章")

                else:
                    messagebox.showwarning("輸入錯誤", "節數錯誤")
                    logging.warning("抓取經文失敗: 節數超出範圍")
        else:
            if verses:
                for v in verses:
                    text_box.insert(tk.END, v + "\n")
                    logging.info(f"成功抓取 {abbr_to_full[book_abbr]} 第 {chapter} 章")

            else:
                text_box.insert(tk.END, "未抓取到經文，請檢查網頁或選擇。")
                logging.error("網頁未回應")
    except Exception as e:
        logging.warning(f"run_search {e}")
        messagebox.showwarning("錯誤", "經文搜索時錯誤")
#創建經文查詢UI
def search_verse_UI():
    try:
        global book_var, chapter_var, verse_var, text_box
        # 標題
        title_label_search = ttk.Label(frame, text="📖 聖經經文查詢", font=("微軟正黑體", 16, "bold"))
        title_label_search.grid(row=0, column=0, columnspan=2, pady=(0, 20))
        # 書卷
        book_label = ttk.Label(frame, text="書卷：", font=("微軟正黑體", 12))
        book_var = tk.StringVar()
        book_combo = ttk.Combobox(frame, textvariable=book_var, values=list(abbr_to_full.keys()), width=15)
        book_label.grid(row=1, column=0, sticky="e", padx=5, pady=5)
        book_combo.grid(row=1, column=1, padx=5, pady=5, sticky="w")

        # 章
        chapter_label = ttk.Label(frame, text="章：", font=("微軟正黑體", 12))
        chapter_var = tk.StringVar()
        chapter_entry = ttk.Entry(frame, textvariable=chapter_var, width=18)
        chapter_label.grid(row=2, column=0, sticky="e", padx=5, pady=5)
        chapter_entry.grid(row=2, column=1, padx=5, pady=5, sticky="w")

        # 節（新加的）
        verse_label = ttk.Label(frame, text="節：", font=("微軟正黑體", 12))
        verse_var = tk.StringVar()
        verse_entry = ttk.Entry(frame, textvariable=verse_var, width=18)
        verse_label.grid(row=3, column=0, sticky="e", padx=5, pady=5)
        verse_entry.grid(row=3, column=1, padx=5, pady=5, sticky="w")

        # 查詢按鈕
        search_btn = ttk.Button(frame, text="查詢", command=run_search)
        search_btn.grid(row=4, column=0, columnspan=2, pady=(15, 0))

        text_box = tk.Text(frame, wrap="word")
        text_box.grid(row=6, column=0, columnspan=2, sticky="nsew", padx=10, pady=10)

        # 置中設定
        for i in range(6):
            frame.grid_rowconfigure(i, weight=1)
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_columnconfigure(1, weight=1)
    except Exception as e:
        logging.warning(f"search_verse_UI {e}")
        messagebox.showwarning("錯誤", "經文搜尋UI錯誤")
#創建PPT的UI 
def produce_the_slide_UI():
    try:
        global progress_bar, status_var
        ttk.Label(frame, text="Word 輸入:").grid(row=0, column=0, columnspan=2, pady=(20, 0), sticky="s")
        ttk.Entry(frame, textvariable=word_path_var, width=50, state='readonly').grid(row=1, column=0, padx=5, pady=5, sticky="e")
        ttk.Button(frame, text="選擇 Word", command=select_word_file).grid(row=1, column=1, padx=5, pady=5, sticky="w")

        ttk.Label(frame, text="PPT 輸出:").grid(row=2, column=0, columnspan=2, pady=(20, 0), sticky="s")
        ttk.Entry(frame, textvariable=ppt_save_var, width=50, state='readonly').grid(row=3, column=0, padx=5, pady=5, sticky="e")
        ttk.Button(frame, text="選擇儲存", command=select_save_path).grid(row=3, column=1, padx=5, pady=5, sticky="w")

        # 按鈕
        produce_btn = ttk.Button(frame, text="製作", command=Start_produce)
        produce_btn.grid(row=4, column=0, columnspan=2, pady=(15, 0))

        progress_bar = ttk.Progressbar(frame, orient="horizontal", mode="determinate")
        status_var = tk.StringVar(value="狀態：待命中")
        progress_bar.grid(row=5, column=0, columnspan=2, padx=5, pady=5, sticky="ew")
        ttk.Label(frame, textvariable=status_var).grid(row=6, column=0, columnspan=2, padx=5, pady=5)

        # 置中設定
        for i in range(7):
            frame.grid_rowconfigure(i, weight=1)
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_columnconfigure(1, weight=1)
    except Exception as e:
        logging.warning(f"produce_the_slide_UI {e}")
        messagebox.showwarning("錯誤", "PPT創建的UI錯誤")
#切換頁面
def change_page():
    global search_page
    try:
        clear_frame(frame)
        search_page = not search_page
        if search_page:
            search_verse_UI()
        else:
            produce_the_slide_UI()
    except Exception as e:
        logging.warning(f"change_page {e}")
        messagebox.showwarning("錯誤", "切換錯誤")
#選取word檔案位置
def select_word_file():
    """打開檔案對話框，讓使用者選擇 Word 檔案 (.docx)"""
    # filedialog.askopenfilename() 打開選擇檔案的對話框
    path = filedialog.askopenfilename(
        title="選擇 Word 證道文件",
        defaultextension=".docx", # 預設副檔名
        filetypes=[
            ("Word 檔案", "*.docx"),
            ("所有檔案", "*.*")
        ]
    )
    if path:
        # 如果使用者選擇了檔案，將路徑設定到 StringVar 變數中
        word_path_var.set(path)
        logging.info(f"選取 Word 檔案: {path}")
#選取PPT存檔位置
def select_save_path():
    """讓使用者指定輸出 PPT 檔案名稱 (.pptx)"""
    path = filedialog.asksaveasfilename(
        title="指定輸出 PPT 檔案名稱",
        defaultextension=".pptx",
        filetypes=[("PowerPoint 檔案", "*.pptx"), ("所有檔案", "*.*")],
        initialfile="證道投影片.pptx"
    )
    if path:
        ppt_save_var.set(path)
        logging.info(f"選取 PPT 儲存路徑: {path}")
# --- 核心安全更新函式 ---
def update_ui_safe(func, *args, **kwargs):
    """將 UI 更新操作安全地排隊到主執行緒中執行 (關鍵技術)"""
    # root.after(0, ...) 確保 UI 操作在主執行緒中執行，防止崩潰。
    root.after(0, lambda: func(*args, **kwargs))
# --- 進度條更新邏輯 ---
def update_progress(current_step, total_tasks, message):
    """在背景執行緒中呼叫，安全地更新進度條和狀態"""
    global status_var
    if total_tasks <= 0:
        return
        
    percent = (current_step / total_tasks) * 100
    
    # 透過 update_ui_safe 更新 Tkinter 元件
    update_ui_safe(progress_bar.config, value=percent)
    update_ui_safe(status_var.set, f"進度 {percent:.2f}% : {message}")

root = tk.Tk()
root.title("FaithSlide")
root.geometry("650x650")
word_path_var = tk.StringVar(value="請選擇 Word 文件...")
ppt_save_var = tk.StringVar(value="請選擇輸出 PPT 檔案名稱...")

# 外框
frame = ttk.Frame(root, padding=20)
frame.grid(row=0, column=0, columnspan=2, sticky="nsew")

change_btn = ttk.Button(root, text="切換", command=change_page)
change_btn.grid(row=1, column=0, pady=(15, 0), sticky="e")

quit_btn = ttk.Button(root, text="退出", command=close_driver)
quit_btn.grid(row=1, column=1, pady=(15, 0), sticky="w")

produce_the_slide_UI()

# --- 確保 root 的權重配置 ---
root.grid_rowconfigure(0, weight=1)      # 讓 Button 所在的第二行 (row=1) 能夠擴展
root.grid_rowconfigure(1, weight=0)
root.grid_columnconfigure(0, weight=1)   # 讓第一列能擴展
root.grid_columnconfigure(1, weight=1)   # 讓第二列能擴展 (因為 frame 跨越了兩列)
# ----------------------------

if __name__ == "__main__":
    root.mainloop()